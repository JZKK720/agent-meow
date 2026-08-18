"""Convert plans/010-client-overview-marp.md into a 16:9 PPTX deck.

Each `---` in the marp source becomes a slide. Headings, tables, bullet
lists, blockquotes, and bold spans are translated to native pptx shapes.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


SRC_ZH = Path(r"C:\Users\1\github-pr\agent-meow\plans\010-client-overview-marp.md")
DST_DIR = Path(r"C:\Users\1\github-pr\agent-meow\plans")

# --- MeowCat / ColorFire brand tokens (mirrored from plans/themes/meowcat.css) ---
# Light / warm cream paper theme, matching the marp PDF deck.
BRAND_PRIMARY = RGBColor(0xE8, 0x65, 0x1A)  # --brand-primary  ember orange
BRAND_ACCENT = RGBColor(0xFF, 0xB3, 0x47)  # --brand-accent   warm amber
BRAND_BG_LIGHT = RGBColor(0xFF, 0xFB, 0xF5)  # --brand-bg-light warm cream paper
BRAND_BG_DARK = RGBColor(0x1A, 0x14, 0x10)  # --brand-bg-dark  warm dark
BRAND_BORDER = RGBColor(0xEA, 0xDF, 0xD4)  # 1px cream-on-cream border
BRAND_TEXT = RGBColor(0x1A, 0x14, 0x10)  # --brand-text
BRAND_TEXT_2 = RGBColor(0x6B, 0x5D, 0x4F)  # --brand-text-secondary
BRAND_TEXT_MUTED = RGBColor(0x9A, 0x8B, 0x7A)  # --brand-text-muted
ROW_ALT = RGBColor(0xFD, 0xF6, 0xE8)  # very light cream row alt

# Backwards-compatible aliases (so the existing code below keeps working)
BG_DARK = BRAND_BG_LIGHT  # slide background
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)  # card / header surface
ACCENT = BRAND_PRIMARY  # ember orange
ACCENT2 = BRAND_ACCENT  # warm amber
TEXT_PRIMARY = BRAND_TEXT  # main text
TEXT_MUTED = BRAND_TEXT_MUTED  # muted text
TEXT_DARK = RGBColor(0xFF, 0xFF, 0xFF)  # text on ember header (white on orange)
BORDER = BRAND_BORDER  # cream border


# ---------------------------------------------------------------------------
# Markdown slicing
# ---------------------------------------------------------------------------


def split_slides(md: str) -> list[list[str]]:
    """Split marp markdown into per-slide blocks. The first frontmatter block
    (before the first `---`) is treated as the title slide."""
    blocks: list[list[str]] = [[]]
    for line in md.splitlines():
        if line.strip() == "---":
            blocks.append([])
        else:
            blocks[-1].append(line)
    # Drop empty trailing blocks
    return [b for b in blocks if any(l.strip() for l in b)]


def classify_block(block: list[str]) -> str:
    """Classify a block as title, table, list, paragraph, or hr."""
    nonblank = [l for l in block if l.strip()]
    if not nonblank:
        return "empty"
    if all(l.lstrip().startswith("|") and l.rstrip().endswith("|") for l in nonblank):
        return "table"
    if all(re.match(r"\s*[-*]\s+", l) for l in nonblank):
        return "list"
    if all(l.lstrip().startswith("#") for l in nonblank):
        return "headings"
    return "paragraph"


def parse_inline(text: str) -> list[tuple[str, bool, bool]]:
    """Parse inline markdown into a list of (text, bold, italic) runs.
    Strips backticks but keeps the text content (code styling applied later
    by treating it as bold-ish for emphasis)."""
    # Order matters: bold (**...**) first, then italic (*...*), then code (`...`).
    pattern = re.compile(
        r"\*\*(?P<bold>[^*]+)\*\*"
        r"|\*(?P<italic>[^*]+)\*"
        r"|`(?P<code>[^`]+)`"
    )
    runs: list[tuple[str, bool, bool]] = []
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            runs.append((text[pos : m.start()], False, False))
        if m.group("bold") is not None:
            runs.append((m.group("bold"), True, False))
        elif m.group("italic") is not None:
            runs.append((m.group("italic"), False, True))
        else:
            runs.append((m.group("code"), True, False))
        pos = m.end()
    if pos < len(text):
        runs.append((text[pos:], False, False))
    return runs


def parse_table(block: list[str]) -> tuple[list[str], list[list[str]]]:
    """Parse a pipe-delimited markdown table into (headers, rows).
    Skips the separator line (---|---)."""
    rows: list[list[str]] = []
    for line in block:
        if not line.strip():
            continue
        if not line.lstrip().startswith("|"):
            continue
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        # Drop the markdown alignment row like | --- | :---: |
        if all(re.fullmatch(r":?-+:?", c) for c in cells):
            continue
        rows.append(cells)
    if not rows:
        return [], []
    return rows[0], rows[1:]


# ---------------------------------------------------------------------------
# pptx builders
# ---------------------------------------------------------------------------


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color=TEXT_PRIMARY,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_runs(
    slide,
    left,
    top,
    width,
    height,
    runs,
    *,
    size: int = 16,
    color=TEXT_PRIMARY,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.15,
):
    """Add a textbox whose content is a list of (text, bold, italic) runs."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    first = True
    for text, bold, italic in runs:
        if not text:
            continue
        # Each new line starts a fresh paragraph
        parts = text.split("\n")
        for i, part in enumerate(parts):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            run = p.add_run()
            run.text = part
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return tb


def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_cream_gradient(slide):
    """MarP applies a soft diagonal warm cream gradient. python-pptx doesn't
    expose gradient background fills directly, so we set a single warm cream
    background that matches the mid-tone of the gradient — visually
    indistinguishable to the viewer in practice."""
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    # The 60% stop of the marp gradient is #fff8ec, which gives the warmest
    # feel. Use that for the slide background.
    fill.fore_color.rgb = RGBColor(0xFF, 0xFA, 0xF3)


def add_accent_bar(slide, *, top, height=Cm(0.12), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_footer(slide, slide_num: int, total: int, *, lang: str = "zh"):
    # Bottom accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(0.6), Cm(17.4), Cm(0.18), Cm(0.18))
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    add_textbox(
        slide,
        Cm(0.9),
        Cm(17.25),
        Cm(15),
        Cm(0.5),
        "agent-meow · 010-client-overview",
        size=9,
        color=TEXT_MUTED,
    )
    add_textbox(
        slide,
        Cm(30),
        Cm(17.25),
        Cm(5),
        Cm(0.5),
        f"{slide_num} / {total}",
        size=9,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_table(slide, headers, rows, *, left, top, width, height, col_widths=None):
    cols = len(headers)
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_PRIMARY
        tf = cell.text_frame
        tf.margin_left = Cm(0.1)
        tf.margin_right = Cm(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        run.font.size = Pt(13)
        run.font.bold = True
        # White text on ember header
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Data rows
    for i, row in enumerate(rows, start=1):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            cell.fill.solid()
            # MarP uses rgba(255,255,255,0.82) for normal rows and
            # rgba(255,179,71,0.04) for zebra rows on the cream background.
            # Approximate that with a near-white card + a very subtle warm
            # zebra.
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF6, 0xE3)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf = cell.text_frame
            tf.margin_left = Cm(0.1)
            tf.margin_right = Cm(0.1)
            tf.word_wrap = True
            for k, (text, bold, italic) in enumerate(parse_inline(cell_text)):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = text
                run.font.size = Pt(12)
                run.font.bold = bold
                run.font.italic = italic
                # MarP renders <strong> inside tables in ember orange.
                if bold:
                    run.font.color.rgb = BRAND_PRIMARY
                else:
                    run.font.color.rgb = BRAND_TEXT
    return table_shape


# ---------------------------------------------------------------------------
# Slide rendering
# ---------------------------------------------------------------------------

prs: Presentation | None = None
BLANK = None  # set inside build() to point at the active presentation


def new_slide() -> "Slide":
    slide = prs.slides.add_slide(BLANK)
    add_cream_gradient(slide)
    return slide


def add_title(slide, text: str, *, size=32, top=Cm(1.2), color=TEXT_PRIMARY):
    add_textbox(slide, Cm(1.5), top, Cm(31), Cm(1.5), text, size=size, bold=True, color=color)
    add_accent_bar(slide, top=Cm(top.cm + 1.7), color=BRAND_PRIMARY)


def add_subtitle(slide, text: str, *, top=Cm(3.0), size=18, color=TEXT_MUTED):
    add_textbox(slide, Cm(1.5), top, Cm(31), Cm(1), text, size=size, color=color, italic=True)


def add_paragraph_block(slide, block, *, top, size=18):
    """Render a block of plain paragraph lines (non-list, non-table)."""
    combined = "\n".join(l.rstrip() for l in block if l.strip())
    runs = parse_inline(combined)
    return add_runs(
        slide, Cm(1.5), top, Cm(31), Cm(12), runs, size=size, color=TEXT_PRIMARY, line_spacing=1.3
    )


def add_list_block(slide, block, *, top, size=18):
    tb = slide.shapes.add_textbox(Cm(1.8), top, Cm(30.5), Cm(12))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in block:
        if not line.strip():
            continue
        # Strip leading bullet marker
        m = re.match(r"\s*[-*]\s+(.*)", line)
        if not m:
            continue
        content = m.group(1)
        runs = parse_inline(content)
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        # Bullet glyph
        bullet_run = p.add_run()
        bullet_run.text = "▸ "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = ACCENT
        bullet_run.font.bold = True
        for text, bold, italic in runs:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = TEXT_PRIMARY
    return tb


# ---------------------------------------------------------------------------
# Per-slide dispatch
# ---------------------------------------------------------------------------


def render_title_slide(slide, block, idx, total, *, lang: str = "zh"):
    """First slide: large title + subtitle + meta. Title slide mirrors the
    marp theme: centered block with a vertical ember->amber accent rail on
    the left (no heading underline)."""
    title = ""
    subtitle = ""
    meta_lines: list[str] = []
    for line in block:
        if line.startswith("# "):
            title = line[2:].strip()
        elif line.startswith("## "):
            subtitle = line[3:].strip()
        elif line.startswith("**"):
            meta_lines.append(line.strip())
    # Vertical ember -> amber accent rail on the left (4px wide, 120px tall)
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(8.5), Cm(0.16), Cm(4.7))
    rail.line.fill.background()
    rail.fill.solid()
    rail.fill.fore_color.rgb = BRAND_PRIMARY
    # Soft gradient effect: overlay a second smaller amber bar near the bottom
    rail2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(12.0), Cm(0.16), Cm(1.2))
    rail2.line.fill.background()
    rail2.fill.solid()
    rail2.fill.fore_color.rgb = BRAND_ACCENT

    add_textbox(
        slide, Cm(2.2), Cm(8.8), Cm(28), Cm(3.0), title, size=44, bold=True, color=BRAND_TEXT
    )
    if subtitle:
        add_textbox(
            slide,
            Cm(2.2),
            Cm(11.5),
            Cm(28),
            Cm(1.5),
            subtitle,
            size=24,
            color=BRAND_TEXT_2,
            italic=True,
        )
    # Meta lines (rendered as bullet-style entries)
    meta_top = 14.0
    for line in meta_lines:
        # Strip leading/trailing ** wrappers
        text = line.strip()
        if text.startswith("**") and text.endswith("**"):
            text = text[2:-2]
        runs = parse_inline(text)
        add_runs(
            slide,
            Cm(2.2),
            Cm(meta_top),
            Cm(28),
            Cm(0.7),
            runs,
            size=15,
            color=BRAND_TEXT_2,
            line_spacing=1.2,
        )
        meta_top += 0.7
    add_footer(slide, idx, total, lang=lang)


def render_content_slide(slide, block, idx, total):
    # Find the first H1 as the slide title
    slide_title = None
    body_start = 0
    for i, line in enumerate(block):
        if line.startswith("# "):
            slide_title = line[2:].strip()
            body_start = i + 1
            break
    if slide_title is None:
        slide_title = block[0].strip().lstrip("# ").strip()
        body_start = 1
    add_title(slide, slide_title, size=30, top=Cm(1.0))
    # Slice remaining content into sub-blocks separated by blank lines
    body = block[body_start:]
    sub_blocks: list[list[str]] = [[]]
    for line in body:
        if line.strip():
            sub_blocks[-1].append(line)
        else:
            sub_blocks.append([])
    sub_blocks = [
        b for b in sub_blocks if b
    ]  # Layout: place elements vertically starting at top Cm(3.3)
    cursor = 3.3  # cm
    for sb in sub_blocks:
        kind = classify_block(sb)
        if kind == "table":
            headers, rows = parse_table(sb)
            ncols = len(headers)
            # Estimate height: 0.7cm per row + header
            nrows = len(rows) + 1
            height_cm = 0.7 * nrows + 0.4
            # Cap at remaining slide space
            available_cm = 17 - cursor
            if height_cm > available_cm:
                height_cm = available_cm
            col_widths = None
            if ncols == 2:
                col_widths = [Cm(11), Cm(20)]
            elif ncols == 3:
                col_widths = [Cm(8), Cm(11.5), Cm(11.5)]
            add_table(
                slide,
                headers,
                rows,
                left=Cm(1.5),
                top=Cm(cursor),
                width=Cm(31),
                height=Cm(height_cm),
                col_widths=col_widths,
            )
            cursor += height_cm + 0.4
        elif kind == "list":
            # Estimate height by line count
            n = sum(1 for l in sb if l.strip())
            height_cm = max(1.2, 0.7 * n)
            add_list_block(slide, sb, top=Cm(cursor), size=16)
            cursor += height_cm + 0.2
        elif kind == "headings":
            for line in sb:
                if line.startswith("## "):
                    add_subtitle(slide, line[3:].strip(), top=Cm(cursor), size=18, color=ACCENT2)
                    cursor += 0.9
                elif line.startswith("### "):
                    add_textbox(
                        slide,
                        Cm(1.8),
                        Cm(cursor),
                        Cm(29),
                        Cm(0.7),
                        line[4:].strip(),
                        size=18,
                        bold=True,
                        color=ACCENT,
                    )
                    cursor += 0.8
        else:
            # Paragraph: collapse into one shape, but if multiple lines, keep
            n = sum(1 for l in sb if l.strip())
            if n == 1 and sb[0].strip().startswith("**") and sb[0].strip().endswith("**"):
                # Section lead-in bold line
                txt = sb[0].strip().strip("*")
                add_textbox(
                    slide,
                    Cm(1.5),
                    Cm(cursor),
                    Cm(31),
                    Cm(0.7),
                    txt,
                    size=18,
                    bold=True,
                    color=ACCENT,
                )
                cursor += 0.8
            else:
                height_cm = max(1.0, 0.55 * n + 0.3)
                add_paragraph_block(slide, sb, top=Cm(cursor), size=16)
                cursor += height_cm + 0.1
        if cursor > 15.5:
            break  # safety: don't overflow
    add_footer(slide, idx, total)


# ---------------------------------------------------------------------------
# English translation map
# ---------------------------------------------------------------------------
# Each entry maps a Chinese (or Chinese-mixed) substring from the ZH marp
# source to its English equivalent. Keys are matched as raw substrings, so
# longer/more specific keys must come first in the dict (Python 3.7+ dict
# preserves insertion order). The map is applied in order to the entire
# marp source before slide rendering.

EN_TRANSLATIONS: dict[str, str] = {
    # Title slide
    "agent-meow 语音代理项目概览": "agent-meow Voice Agent Project Overview",
    "灵创K16 + 橘宝R16 双平台方案": "Lingchuang K16 + Jubao R16 Dual-Platform Plan",
    "**目标**：在两台笔记本上实现零云端、零外部 GPU 的本地 AI 语音代理": "**Goal**: Deliver a zero-cloud, zero-external-GPU local AI voice agent on two laptops",
    "**日期**：2026-08-04 · **状态**：灵创K16 方案已就绪，橘宝R16 评估中": "**Date**: 2026-08-04 · **Status**: Lingchuang K16 plan ready, Jubao R16 under evaluation",
    # Slide 2 — problem statement
    "项目目标：解决什么问题？": "Project Goal: What problem are we solving?",
    "**当前痛点**：语音代理冷启动需 **90 秒**，用户体验差": "**Pain point**: Voice agent cold-start takes **90 seconds** — poor user experience",
    "**目标效果**：": "**Target outcome**:",
    "语音启动": "Voice startup",
    "云端依赖": "Cloud dependency",
    "运行成本": "Operating cost",
    "AI 能力": "AI capability",
    "必须联网": "Always online",
    "可选（离线可用）": "Optional (works offline)",
    "API 持续计费": "API billing",
    "零（离线模式）": "Zero (offline mode)",
    "仅对话": "Chat only",
    "对话 + 工具调用（代码/文件）": "Chat + tool calls (code/files)",
    "**核心价值**：语音即入口 — 用户说话即可与 AI 交互，AI 可调用工具完成任务": "**Core value**: Voice is the entry point — users speak to interact with AI, AI invokes tools to complete tasks",
    # Slide 3 — device comparison
    "两台设备一览": "Two Devices at a Glance",
    "灵创K16": "Lingchuang K16",
    "橘宝R16": "Jubao R16",
    "**定位**": "**Positioning**",
    "高端旗舰，大模型质量优先": "Premium flagship, max model quality",
    "性价比之选，CUDA 加速优先": "Best value, CUDA acceleration priority",
    "**处理器**": "**CPU**",
    "**显卡**": "**GPU**",
    "内置 Radeon 8060S (96GB)": "Built-in Radeon 8060S (96GB)",
    "内置 Radeon 890M + **RTX 5060 (8GB)**": "Built-in Radeon 890M + **RTX 5060 (8GB)**",
    "**AI 芯片**": "**AI chip**",
    "**内存**": "**Memory**",
    "**优势**": "**Strength**",
    "96GB 显存跑最大模型": "96GB VRAM runs the largest models",
    "RTX 5060 CUDA 原生加速": "RTX 5060 native CUDA acceleration",
    "**同一软件，两种硬件**：agent-meow 在两台设备上提供相同功能，性能特征不同": "**Same software, two hardware targets**: agent-meow delivers the same functionality on both devices, with different performance profiles",
    # Slide 4 — phased plan
    "方案概述：4 个阶段": "Plan Overview: 4 Phases",
    "**1**": "**1**",
    "**2**": "**2**",
    "**3**": "**3**",
    "**4**": "**4**",
    "云端语音 + GPU 语音识别": "Cloud voice + GPU speech recognition",
    "用户界面集成 + 混合模式": "UI integration + hybrid mode",
    "AI 工具调用能力": "AI tool-calling capability",
    "本地大模型推理": "Local LLM inference",
    "并行开发，互不依赖": "Parallel development, no inter-dependency",
    "依赖阶段 1": "Depends on Phase 1",
    "依赖阶段 3": "Depends on Phase 3",
    "**总工作量**：5 个计划 (006–010)，预计 2–4 周": "**Total effort**: 5 plans (006–010), estimated 2–4 weeks",
    "阶段": "Phase",
    "内容": "Scope",
    "说明": "Notes",
    # Slide 5 — Phase 1
    "阶段 1：云端语音 + GPU 加速": "Phase 1: Cloud Voice + GPU Acceleration",
    "**计划 006**：安装语音网关 + 阿里云 DashScope": "**Plan 006**: Install voice gateway + Alibaba DashScope",
    "解决 90 秒冷启动 → **即时响应**（云端常驻）": "Solve 90-second cold start → **instant response** (cloud always-on)",
    "中文语音质量优秀，中国可直连": "Excellent Mandarin voice quality, direct connection from China",
    "90 天免费试用，之后 ~¥0.20/分钟": "90-day free trial, then ~¥0.20/minute",
    "**计划 008**：GPU 语音识别加速": "**Plan 008**: GPU speech recognition acceleration",
    "灵创K16：用 Vulkan 技术将语音识别放到 GPU（60s → ~3s）": "Lingchuang K16: Move speech recognition onto the GPU with Vulkan (60s → ~3s)",
    "橘宝R16：RTX 5060 原生支持 CUDA，**直接安装即可**（60s → ~1s）": "Jubao R16: RTX 5060 supports CUDA natively — **drop-in install** (60s → ~1s)",
    "**结果**：在线模式即时响应，离线模式 3–8 秒启动": "**Result**: Instant response in online mode, 3–8 second start in offline mode",
    # Slide 6 — Phase 2
    "阶段 2：用户界面 + 混合模式": "Phase 2: User Interface + Hybrid Mode",
    "**计划 006b**：在线/离线自动切换": "**Plan 006b**: Online/offline auto-switching",
    "云端不可用时自动切离线，用户无感知": "Auto-fall-back to offline when cloud is unavailable — no user-visible change",
    "解决中国网络环境下的可用性问题": "Solves availability issues in mainland-China network conditions",
    "**计划 007**：保留猫爪 UI，替换底层引擎": "**Plan 007**: Keep the MeowCat paw UI, swap the underlying engine",
    "用户看到的界面不变（猫爪按钮 + 波形动画）": "What the user sees stays the same (paw button + waveform animation)",
    "内部传输层升级为更成熟的 QAA 协议": "The internal transport layer upgrades to the more mature QAA protocol",
    "新增在线/离线切换按钮": "Add an online/offline toggle button",
    "**结果**：用户体验一致，底层更稳定": "**Result**: Consistent user experience, more stable under the hood",
    # Slide 7 — Phase 3+4
    "阶段 3+4：AI 工具调用 + 本地大模型": "Phase 3+4: AI Tool Calling + Local LLM",
    "**计划 009**：AI 语音 → 工具调用": "**Plan 009**: AI voice → tool calls",
    "用户说话 → 简单问题即时回答": "User speaks → simple questions get an instant answer",
    "需要工具时（代码/文件/搜索）→ 后台执行 → 语音播报结果": "When tools are needed (code/files/search) → run in background → speak the result",
    "AI 语音成为工作入口，不仅仅是聊天": "AI voice becomes a work interface, not just a chat surface",
    "**计划 010**：本地大模型推理": "**Plan 010**: Local LLM inference",
    "灵创K16：35B 大模型完全驻留 96GB 显存（最高质量）": "Lingchuang K16: 35B model lives entirely in 96GB VRAM (highest quality)",
    "橘宝R16：同一 35B 模型，智能分配到 RTX 5060 + 系统内存": "Jubao R16: Same 35B model, intelligently split between RTX 5060 + system RAM",
    "**零 API 费用**，完全本地运行": "**Zero API fees**, fully local execution",
    "**结果**：语音驱动的 AI 代理，可执行任务，零云端成本": "**Result**: A voice-driven AI agent that executes tasks with zero cloud cost",
    # Slide 8 — final comparison
    "两台设备最终效果对比": "Final Outcome Comparison",
    "指标": "Metric",
    "语音识别": "Speech recognition",
    "AI 模型": "AI model",
    "35B 大模型（最高质量）": "35B model (highest quality)",
    "同 35B 模型（稍压缩）": "Same 35B model (slightly compressed)",
    "离线运行": "Offline operation",
    "✅ 零云端": "✅ Zero cloud",
    "实现难度": "Implementation effort",
    "中等（需编译）": "Medium (compile required)",
    "**简单**（直接安装）": "**Simple** (drop-in install)",
    "**总结**：灵创K16 以大显存追求最高质量；橘宝R16 以 CUDA 追求最快速度和最简部署": "**Summary**: Lingchuang K16 chases the highest quality with its huge VRAM; Jubao R16 chases the fastest speed and simplest setup with CUDA",
    # Slide 9 — cost & risk
    "成本与风险": "Cost & Risk",
    "**开发成本**：": "**Development cost**:",
    "项目": "Item",
    "云端语音": "Cloud voice",
    "90天免费，后 ~¥0.20/分": "90-day free, then ~¥0.20/min",
    "同": "Same",
    "离线运行": "Offline operation",
    "**零成本**": "**Zero cost**",
    "开发工作量": "Dev effort",
    "中等": "Medium",
    "**较低**": "**Lower**",
    "**风险**：": "**Risk**:",
    "风险": "Risk",
    "级别": "Level",
    "缓解措施": "Mitigation",
    "云端服务不可用": "Cloud service outage",
    "低": "Low",
    "中": "Medium",
    "自动切离线模式": "Auto-switch to offline mode",
    "GPU 兼容性": "GPU compatibility",
    "已验证两平台驱动": "Drivers verified on both platforms",
    "模型性能不足": "Model performance insufficient",
    "可切换更小模型": "Can switch to a smaller model",
    # Slide 10 — vision
    "项目愿景": "Project Vision",
    "**agent-meow** 是一个语音优先的 AI 代理，用户通过说话即可：": "**agent-meow** is a voice-first AI agent — users speak to:",
    "💬 **对话**：自然语言交流，中英文双语": "💬 **Chat**: Natural language conversation, bilingual EN/ZH",
    "🔧 **执行任务**：调用代码、文件、搜索等工具": "🔧 **Execute tasks**: Invoke code, files, search, and other tools",
    "🏠 **完全本地**：离线可用，零云端费用": "🏠 **Fully local**: Works offline, zero cloud fees",
    "🎨 **品牌体验**：猫爪 UI，独特的交互设计": "🎨 **Brand experience**: MeowCat paw UI, distinctive interaction design",
    "**双平台覆盖**：": "**Dual-platform coverage**:",
    "灵创K16 面向高端用户（96GB 显存，最大模型）": "Lingchuang K16 for power users (96GB VRAM, max-size models)",
    "橘宝R16 面向主流用户（RTX 5060 CUDA，性价比高）": "Jubao R16 for mainstream users (RTX 5060 CUDA, best price/perf)",
    "两个平台运行同一软件，提供同一体验，适配不同硬件。": "Same software on both platforms, same user experience, tuned to different hardware.",
    "当前": "Current",
    "目标": "Target",
    "**90 秒**": "**90 seconds**",
    "~8 秒": "~8 seconds",
    "~3 秒": "~3 seconds",
    "**即时**（在线）/ ~8 秒（离线）": "**Instant** (online) / ~8 seconds (offline)",
    "**即时**（在线）/ ~8s（离线）": "**Instant** (online) / ~8s (offline)",
    "**即时**（在线）/ ~3s（离线）": "**Instant** (online) / ~3s (offline)",
    "**即时**": "**Instant**",
    "**可选**（离线可用）": "**Optional** (works offline)",
    "**零**（离线模式）": "**Zero** (offline mode)",
    "**对话 + 工具调用**（代码/文件）": "**Chat + tool calls** (code/files)",
    "**可选**": "**Optional**",
    "**零**": "**Zero**",
    "**对话 + 工具调用**": "**Chat + tool calls**",
    "秒（": "second (",
    # Single-character / short tokens that appear inside tables and need
    # their own entries (these won't match any of the longer keys above)
    "云端": "Cloud",
    "语音": "Voice",
    "阶段 ": "Phase ",
    "计划": "Plan",
    "橘宝": "Jubao",
    "灵创": "Lingchuang",
}


def translate_to_en(zh_md: str) -> str:
    """Apply the EN_TRANSLATIONS map to a ZH marp source, longest keys first."""
    out = zh_md
    # Sort by key length descending so longer substrings replace before
    # any shorter prefix substring (e.g. "灵创K16" before "灵创").
    for zh, en in sorted(EN_TRANSLATIONS.items(), key=lambda kv: -len(kv[0])):
        if zh in out:
            out = out.replace(zh, en)
    return out


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build(lang: str) -> Path:
    """Build a single PPTX deck. `lang` is 'zh' or 'en'."""
    global prs, BLANK
    zh_md = SRC_ZH.read_text(encoding="utf-8")
    if lang == "zh":
        md = zh_md
    elif lang == "en":
        md = translate_to_en(zh_md)
    else:
        raise ValueError(f"Unknown lang: {lang}")

    # Re-init a fresh Presentation per language so the previous run's slides
    # don't leak into the next build.
    prs = Presentation()
    prs.slide_width = Cm(33.867)  # 16:9 @ 13.33"
    prs.slide_height = Cm(19.05)
    BLANK = prs.slide_layouts[6]

    slides = split_slides(md)
    # The first marp block in the source is the frontmatter / title slide.
    # The second block is also title/subtitle context, not a content slide.
    # We treat block 0 (the H1 + ## + meta lines) as the title slide and skip
    # block 1 because it only repeats the title pair.
    total = len(slides)
    rendered = 0
    for i, block in enumerate(slides):
        # Skip the redundant second block in this marp file (only contains the
        # ## subtitle line that already appears as the H2 on the title slide).
        if i == 1 and len(block) <= 2 and all(l.startswith("#") for l in block if l.strip()):
            continue
        slide = new_slide()
        rendered += 1
        if i == 0:
            render_title_slide(slide, block, rendered, total - 1, lang=lang)
        else:
            render_content_slide(slide, block, rendered, total - 1)
    out = DST_DIR / f"010-client-overview-{lang}.pptx"
    prs.save(str(out))
    print(f"Wrote {out} ({rendered} slides, lang={lang})")
    return out


def main() -> None:
    build("zh")
    build("en")


if __name__ == "__main__":
    main()
